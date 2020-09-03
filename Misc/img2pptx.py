# A nifty script to convert images in a folder to a ppt.
# Built by @radiantly for KK

# Installing dependencies:
# pip install python-pptx natsort tqdm

# Usage:
# python img2pptx.py path_to_images_folder
# Tested on Linux / Python 3.8.5

import sys

import pptx
import pptx.util
from pptx.enum.shapes import MSO_SHAPE

from tqdm import tqdm
from pathlib import Path
from natsort import natsorted, ns


def makePresentation(images: list, saveDirectory: Path):
    prs = pptx.Presentation()

    # slide height @ 4:3
    # prs.slide_height = 6858000
    # slide height @ 16:9
    # prs.slide_height = 5143500

    # print(prs.slide_width, prs.slide_height)

    for image in tqdm(images):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        pic = slide.shapes.add_picture(image.as_posix(), shape.left, shape.top)

        # calculate max width/height for target size
        ratio = min(shape.width / float(pic.width), shape.height / float(pic.height))

        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)

        pic.left = shape.left + ((shape.width - pic.width) // 2)
        pic.top = shape.top + ((shape.height - pic.height) // 2)

        placeholder = shape.element
        placeholder.getparent().remove(placeholder)

    pptxSaveLocation = saveDirectory / "images.pptx"
    prs.save(pptxSaveLocation)
    print(f"PPt successfully saved to {pptxSaveLocation.as_posix()}")


def main():
    if len(sys.argv) != 2:
        from tkinter import Tk, filedialog

        Tk().withdraw()
        print("Select folder with images ..")
        imagesFolderPath = Path(filedialog.askdirectory())
    else:
        imagesFolderPath = Path(sys.argv[1])

    if not imagesFolderPath.exists():
        raise Exception(f"Folder {imagesFolderPath.as_posix()} does not exist")

    # Retrieve images
    images = [
        image for image in imagesFolderPath.glob("*") if image.suffix in [".jpeg", ".jpg", ".png"]
    ]

    if not len(images):
        raise Exception(f"Cannot find any images in {imagesFolderPath.as_posix()}")

    print(f"Found {len(images)} images")

    # Sort images in order
    imagesSorted = natsorted(images, alg=ns.PATH)

    makePresentation(imagesSorted, imagesFolderPath)


if __name__ == "__main__":
    main()
